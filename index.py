# -*- coding: utf-8 -*-

import arabic_reshaper

import pptx
from pptx import Presentation , text
from pptx.enum.shapes import MSO_SHAPE
import os
from pptx.util import Inches, Pt
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt , Cm
from pptx.enum.text import PP_ALIGN
from PIL import Image




class creat_video():
    def __init__(self):
        self.p = open('hemidi.pptx', 'rb')
        self.prs = Presentation(self.p)

        #########################  ketba
        #shape = prs.slides[0]
        ###############
        #print len(self.prs.slides)
        self.slide_0 = self.prs.slides[0]
        self.slide_1 = self.prs.slides[1]
        self.slide_2 = self.prs.slides[2]
        self.slide_3 = self.prs.slides[3]
        self.slide_4 = self.prs.slides[4]
        self.slide_5 = self.prs.slides[5]
        self.slide_6 = self.prs.slides[6]
        self.slide_7 = self.prs.slides[7]

        #shapes1 = slide.shapes
        #print len(shapes1)

        #print type(shapes1)

    ##### cett fonction pour chenge photo as_ map_ embdad
    def image_(self,patch_img , index_shape):

        # slide0_1
        #if slide_0_0 == None :
        REPLACEMENT_IMG = str(patch_img)
        img = self.slide_0.shapes[index_shape]
        # get current image info
        imgPic = img._pic
        imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
        imgPart = self.slide_0.part.related_parts[imgRID]

        # get info about replacement image
        with open(REPLACEMENT_IMG, 'rb') as f:
            rImgBlob = f.read()
        rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
        rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

        # replace
        imgPart._blob = rImgBlob

        # now alter the size and position to suit that of the replacement img
        # rescale sizes so image isn't stretched
        widthScale = float(rImgWidth) / img.width
        heightScale = float(rImgHeight) / img.height
        maxScale = max(widthScale, heightScale)
        # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
        # center the image if it's different size to the original
        if index_shape in [4,6,10]:
            scaledImgWidth, scaledImgHeight = Cm(14.75), Cm(7.63)
            scaledImgLeft = Cm(2.16)  # int(img.left + (img.width - scaledImgWidth) / 2)
            scaledImgTop = Cm(3.79)  # int(img.top + (img.height - scaledImgHeight) / 2)

        elif index_shape in [3,5,9]:
            scaledImgWidth, scaledImgHeight = Cm(14.75), Cm(7.63)
            scaledImgLeft = Cm(17.49)  # int(img.left + (img.width - scaledImgWidth) / 2)
            scaledImgTop = Cm(3.77)  # int(img.top + (img.height - scaledImgHeight) / 2)

        # now update
        img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight







    #### ceet fonction  pour  suprémme les  charachter  *1 ... à les text dant  shape
    def delit_smile(self,x) :
        l = ['0', '1', '2', '3', '4', '5', '6', '7', '8','9','*','#']
        output = ''
        for  i  in x :
            if i  in l :
                pass
            else :
                output = output + i
        return  output
    ##### fine  les  fonction

    ########################################## format nomre  #####
    def format_(self,x):
        """
        for i in str(x):
            if i in '0123456789':
                return  str("{:,}".format(int(x)))
            else:
                return x
        """
        i = 0
        k = ''
        #a = ', '.join(x.encode('utf-8')
        if type(x) == int :
            x = str(x)

        while i < len(x):
            #x[i] = ','.join(str(x[i]))
            l = x[i]


            if unicode(l) not in '0123456789.,':
                break
            i = i + 1
        if i == len(x):
            l = x
            try:
                return str("{:,}".format(int(unicode(l))))
            except :
                return  str(x)
        else:
            return x


    ############################### fir  f m

    ####### cett fonction for  pour ajouté  les nouveau  text  dans  chape
    def ajouter_text(self,valeur_stat_slid0_1, valeur_stat_slid0_2 , valeur_stat_slid1_1,
                           valeur_stat_slid1_2,valeur_stat_slid2_1,valeur_stat_slid2_2,
                           valeur_stat_slid3_1,valeur_stat_slid3_2,slide_0_0,slide_1_1,image,valeur_stat_slid4_1,valeur_stat_slid4_2,
                     valeur_stat_slid5_1,valeur_stat_slid5_2,image_as,image_embdad,image_map,valeur_stat_slid6_1 , valeur_stat_slid6_2,
                     valeur_stat_slid7_1,valeur_stat_slid7_2) :

        ###################################### slide 0 ####################################
        #slid0_1
        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        #valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None :
            for shape in self.slide_0.shapes:
                for i in smile_slide0_1 :
                    smaile = i
                    val = self.format_(valeur_stat_slid0_1[smile_slide0_1.index(smaile)])
                    #val = valeur_stat_slid0_1[smile_slide0_1.index(smaile)]
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if  smaile in text_frame.text :
                        disc = self.delit_smile(text_frame.text)
                        #val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text =  '%s  %s'%(disc,val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        #print 'done'

            ################## chengi les  photo  #########################
            # image flag slid0_ 1 ===> 21
            # image flag slid0_ 2 ===> 29
            # image petirt flag slid0_ 1 ===> 37
            # image petit flag slid0_ 2 ===> 38
            for i in [1,41] :

                #slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_0.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_0.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth)  , Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                #scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                if  i ==  1 :
                    scaledImgWidth, scaledImgHeight = Cm(14.42), Cm(7.66)
                    scaledImgLeft = Cm(2.16) #int(img.left + (img.width - scaledImgWidth) / 2)
                    scaledImgTop =  Cm(3.77) #int(img.top + (img.height - scaledImgHeight) / 2)

                else :
                    scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                    scaledImgLeft =  int(img.left + (img.width - scaledImgWidth) / 2)
                    scaledImgTop =  int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
            for i  in  [4,6,10] :
                if i == 4 :
                    self.image_(str(image_map),4)
                elif i == 6 :
                    self.image_(str(image_embdad), 6)
                elif i == 10 :
                    self.image_(str(image_as), 10)


        #slide0.2
        if slide_1_1 == None :
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            #valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_0.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    #str("{:,}".format(int(self.line1.text())))
                    #val_2 = valeur_stat_slid0_2[smile_slide0_1.index(smaile)]
                    val_2 = self.format_(valeur_stat_slid0_2[smile_slide0_2.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        #print 'done'
            ################## chengi les  photo  #########################
            # image flag slid0_ 1 ===> 21
            # image flag slid0_ 2 ===> 29
            # image petirt flag slid0_ 1 ===> 37
            # image petit flag slid0_ 2 ===> 38
            # image petit flag slid1_ 1 ===> 41
            # image petit flag slid1_ 2 ===> 42
            # image petit flag slid2_ 1 ===> 41
            # image petit flag slid2_ 2 ===> 42
            # image petit flag slid3_ 1 ===> 40
            # image petit flag slid3_ 2 ===> 41
            for i in [2, 42]:


                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_0.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_0.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                if i == 2:
                    scaledImgWidth, scaledImgHeight = Cm(14.75), Cm(7.63)
                    scaledImgLeft = Cm(17.49)  # int(img.left + (img.width - scaledImgWidth) / 2)
                    scaledImgTop = Cm(3.79)  # int(img.top + (img.height - scaledImgHeight) / 2)

                else:
                    scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                    scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                    scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight


            for i  in  [3,5,9] :
                if i == 3 :
                    self.image_(str(image_map),3)
                elif i == 5 :
                    self.image_(str(image_embdad), 5)
                elif i == 9 :
                    self.image_(str(image_as), 9)


        ##################### save

        ################################### slide 1  #########################

        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        # valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None:
            for shape in self.slide_1.shapes:
                for i in smile_slide0_1:
                    smaile = i
                    #val = valeur_stat_slid1_1[smile_slide0_1.index(smaile)]
                    val = self.format_(valeur_stat_slid1_1[smile_slide0_1.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        # print 'done'
            #################### add petit emage from  slid1 ####################
            for i in [41]:


                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_1.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_1.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight



        # slide0.2
        if slide_1_1 == None:
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            # valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_1.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    #val_2 = valeur_stat_slid1_2[smile_slide0_2.index(smaile)]
                    val_2 =self.format_(valeur_stat_slid1_2[smile_slide0_2.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
            #################### add petit emage from  slid1 ####################
            for i in [42]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_1.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_1.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight

        ##########################################  slide 2 ##########################################################

        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        # valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None:
            for shape in self.slide_2.shapes:
                for i in smile_slide0_1:
                    smaile = i
                    #val = valeur_stat_slid2_1[smile_slide0_1.index(smaile)]
                    val = self.format_(valeur_stat_slid2_1[smile_slide0_1.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        # print 'done'
            #################### add petit emage from  thise slid1 ####################
            for i in [41]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_2.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_2.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight



        # slide0.2
        if slide_1_1 == None:
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            # valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_2.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    #val_2 = valeur_stat_slid2_2[smile_slide0_2.index(smaile)]
                    val_2 = self.format_(valeur_stat_slid2_2[smile_slide0_2.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0,0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
            #################### add petit emage from  thise slid1 ####################
            for i in [42]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_2.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_2.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight




        ############################# slide 3 ###################################

        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        # valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None:
            for shape in self.slide_3.shapes:
                for i in smile_slide0_1:
                    #print 'benameur'
                    smaile = i
                    #val = valeur_stat_slid3_1[smile_slide0_1.index(smaile)]
                    val = self.format_(valeur_stat_slid3_1[smile_slide0_1.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0,0,0)
                        p.alignment = PP_ALIGN.CENTER
                        # print 'done'

            #################### add petit emage from  thise slid1 ####################
            for i in [41]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_3.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_3.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight


        # slide0.2

        if slide_1_1 == None:
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            # valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_3.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    #val_2 = valeur_stat_slid3_2[smile_slide0_2.index(smaile)]
                    val_2 = self.format_(valeur_stat_slid3_2[smile_slide0_2.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER

            #################### add petit emage from  thise slid1 ####################
            for i in [42]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_3.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_3.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight


        ############################################################################################################
        ############################################################################################################
        ############################################################################################################
        ############################################################################################################
        ############################################################################################################
        ############################## SLIDE 4 5 6 #################################################################
        ############################################################################################################
        ############################################################################################################
        ############################################################################################################
        ############################################################################################################
        ############################################################################################################
        ############################################################################################################
        ############################################################################################################

        ##################################### slide 4 ########################################################################

        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        #valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None:
            for shape in self.slide_4.shapes:
                for i in smile_slide0_1:
                    #print 'hemidi'
                    smaile = i
                    # val = valeur_stat_slid2_1[smile_slide0_1.index(smaile)]
                    val = self.format_(valeur_stat_slid4_1[smile_slide0_1.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        # print 'done'
            #################### add petit emage from  thise slid1 ####################

            for i in [41]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_4.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_4.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight

        # slide0.2

        if slide_1_1 == None:
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            # valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_4.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    # val_2 = valeur_stat_slid2_2[smile_slide0_2.index(smaile)]
                    val_2 = self.format_(valeur_stat_slid4_2[smile_slide0_2.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
            #################### add petit emage from  thise slid1 ####################


            for i in [42]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_4.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_4.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
        ############################# slide 5 ###################################


        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        # valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None:
            for shape in self.slide_5.shapes:
                for i in smile_slide0_1:
                    smaile = i
                    # val = valeur_stat_slid3_1[smile_slide0_1.index(smaile)]
                    val = self.format_(valeur_stat_slid5_1[smile_slide0_1.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        # print 'done'

            #################### add petit emage from  thise slid1 ####################
            for i in [43]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_5.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_5.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight

        # slide0.2
        if slide_1_1 == None:
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            # valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_5.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    # val_2 = valeur_stat_slid3_2[smile_slide0_2.index(smaile)]
                    val_2 = self.format_(valeur_stat_slid5_2[smile_slide0_2.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER

            #################### add petit emage from  thise slid1 ####################
            for i in [42]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_5.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_5.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
        ############################# slide 6 ###################################

        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        # valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None:
            print  'hemidi benameur'
            for shape in self.slide_6.shapes:
                for i in smile_slide0_1:
                    smaile = i
                    # val = valeur_stat_slid3_1[smile_slide0_1.index(smaile)]
                    val = self.format_(valeur_stat_slid6_1[smile_slide0_1.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        print  'Yessssssss1111111111111111111'
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        # print 'done'

            #################### add petit emage from  thise slid1 ####################
            """
            for i in [43]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_6.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_6.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
                """

        # slide0.2
        if slide_1_1 == None:
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            # valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_6.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    # val_2 = valeur_stat_slid3_2[smile_slide0_2.index(smaile)]
                    val_2 = self.format_(valeur_stat_slid6_2[smile_slide0_2.index(smaile)])

                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        print("YESSSS222222222222222")
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER

            #################### add petit emage from  thise slid1 ####################
            """
            for i in [42]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_6.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_6.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
                """

        ############################# slide 7 ###################################

        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        # valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None:
            for shape in self.slide_7.shapes:
                for i in smile_slide0_1:
                    smaile = i
                    # val = valeur_stat_slid3_1[smile_slide0_1.index(smaile)]
                    val = self.format_(valeur_stat_slid7_1[smile_slide0_1.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:

                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        # print 'done'

            #################### add petit emage from  thise slid1 ####################
            """
            for i in [43]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_7.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_7.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
                """

        # slide0.2
        if slide_1_1 == None:
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            # valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_7.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    # val_2 = valeur_stat_slid3_2[smile_slide0_2.index(smaile)]
                    val_2 = self.format_(valeur_stat_slid7_2[smile_slide0_2.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER

            #################### add petit emage from  thise slid1 ####################
            """
            for i in [42]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_7.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_7.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
                """

        '''

        ############################# slide 8 ###################################

        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        # valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None:
            for shape in self.slide_8.shapes:
                for i in smile_slide0_1:
                    smaile = i
                    # val = valeur_stat_slid3_1[smile_slide0_1.index(smaile)]
                    val = self.format_(f1[smile_slide0_1.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        # print 'done'

            #################### add petit emage from  thise slid1 ####################
            for i in [43]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_8.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_8.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight

        # slide0.2
        if slide_1_1 == None:
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            # valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_8.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    # val_2 = valeur_stat_slid3_2[smile_slide0_2.index(smaile)]
                    val_2 = self.format_(f1[smile_slide0_2.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER

            #################### add petit emage from  thise slid1 ####################
            for i in [42]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_8.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_8.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight

            ############################# slide 9 ###################################

        smile_slide0_1 = ['0*', '1*', '2*', '3*', '4*', '5*', '6*', '7*', '8*', '9*']
        # valeur_stat_slid0_1 = [100, 200, 300, 400, 500, 600, 700, 800, 900, 'hemidi']
        if slide_0_0 == None:
            for shape in self.slide_9.shapes:
                for i in smile_slide0_1:
                    smaile = i
                    # val = valeur_stat_slid3_1[smile_slide0_1.index(smaile)]
                    val = self.format_(f1[smile_slide0_1.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER
                        # print 'done'

            #################### add petit emage from  thise slid1 ####################
            for i in [43]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_9.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_9.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight

        # slide0.2
        if slide_1_1 == None:
            smile_slide0_2 = ['0#', '1#', '2#', '3#', '4#', '5#', '6#', '7#', '8#', '9#']
            # valeur_stat_slid0_2 = [900, 800, 700, 600, 500, 400, 300, 200, 100, 'hemidi']
            for shape in self.slide_9.shapes:
                for i in smile_slide0_2:
                    smaile = i
                    # val_2 = valeur_stat_slid3_2[smile_slide0_2.index(smaile)]
                    val_2 = self.format_(f1[smile_slide0_2.index(smaile)])
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        disc = self.delit_smile(text_frame.text)
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s  %s' % (disc, val_2)
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER

            #################### add petit emage from  thise slid1 ####################
            for i in [42]:
                # slide = presso.slides[0]
                REPLACEMENT_IMG = str(image)
                img = self.slide_9.shapes[i]

                # get current image info
                imgPic = img._pic
                imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                imgPart = self.slide_9.part.related_parts[imgRID]

                # get info about replacement image
                with open(REPLACEMENT_IMG, 'rb') as f:
                    rImgBlob = f.read()
                rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                rImgWidth, rImgHeight = Cm(rImgWidth), Cm(rImgHeight)  # change from Px

                # replace
                imgPart._blob = rImgBlob

                # now alter the size and position to suit that of the replacement img
                # rescale sizes so image isn't stretched
                widthScale = float(rImgWidth) / img.width
                heightScale = float(rImgHeight) / img.height
                maxScale = max(widthScale, heightScale)
                # scaledImgWidth, scaledImgHeight = Cm(14.42),Cm(7.66) #int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                # center the image if it's different size to the original
                scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                # now update
                img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
                '''

        self.prs.save("hemidi.pptx")



    def save_fiile(self):
        self.prs.save("hemidi.pptx")
        self.p.close()



"""save the file"""
#prs.save("filezip.pptx")
#ajouter_text()
#m = creat_video()
#m.ajouter_text()