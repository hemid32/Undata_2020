#add_shape(autoshape_type_id, left, top, width, height)[source]
# -*- coding: utf-8 -*-

import arabic_reshaper
from pickle import dump, load

import pptx
from pptx import Presentation , text
from pptx.enum.shapes import MSO_SHAPE
import os
from pptx.util import Inches, Pt
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE


from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt , Cm
from pptx.enum.text import PP_ALIGN
from PIL import Image
import sys
reload(sys)
sys.setdefaultencoding('utf-8')

def calcul(x) :
    try:
        f = open('stor.bin', 'r')
        l = open('stor2.bin', 'r')
        m2 = load(l)
        m1 = load(f)
        REPLACEMENT_IMG = m1[-1]
        REPLACEMENT_IMG2 = m2[-1]
        f.close()
        l.close()

        p_ = open('xc.pptx', 'rb')
        prs = Presentation(p_)
        print x
        u = 0
        for cnt in x :
            h = prs.slides[u]
            #h = prs.slides[2]
            u = u+1
            #contry 1
            contry1 = round(cnt[0]/float(10),2)
            contry2 = round(cnt[1] / float(10), 2)
            h.shapes.add_shape(MSO_SHAPE.RECTANGLE,Cm(8.18),Cm(8.81),Cm(2.14 ),Cm(contry2))
            #contry 2

            h.shapes.add_shape(MSO_SHAPE.RECTANGLE,Cm(23.46),Cm(8.81),Cm(2.14 ),Cm(contry1))

            for shape in h.shapes :
                l = shape
                if l.shape_id in [62,63,40,41] :
                    l.fill.solid()
                    l.line.fill.solid()
                    l.fill.fore_color.rgb = RGBColor(242, 242, 242)
                    l.line.fill.fore_color.rgb = RGBColor(242, 242, 242)

                smile_slide0_1 = ['*','#']
                for i in smile_slide0_1:
                    smaile = i
                    if i == '*' :
                        val = cnt[0]
                    elif i == '#':
                        val = cnt[1]
                    # val = valeur_stat_slid0_1[smile_slide0_1.index(smaile)]
                    if not shape.has_text_frame:
                        continue
                    text_frame = shape.text_frame

                    if smaile in text_frame.text:
                        # val = 55555
                        text_frame.text = ''
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = '%s'%(str(val)+' %')
                        font = run.font
                        font.name = 'Adobe Gothic Std B'
                        font.size = Pt(38)
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.rgb = RGBColor(0, 0, 0)
                        font.fill.solid()
                        font.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        p.alignment = PP_ALIGN.CENTER

                if l.shape_id == 7 :
                    img = h.shapes[7]  # .element[0]
                    imgPic = img._pic
                    imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                    imgPart = h.part.related_parts[imgRID]

                    # get info about replacement image
                    with open(REPLACEMENT_IMG, 'rb') as f:
                        rImgBlob = f.read()
                    rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                    rImgWidth, rImgHeight = Pt(rImgWidth), Pt(rImgHeight)  # change from Px

                    # replace
                    imgPart._blob = rImgBlob

                    # now alter the size and position to suit that of the replacement img
                    # rescale sizes so image isn't stretched
                    widthScale = float(rImgWidth) / img.width
                    heightScale = float(rImgHeight) / img.height
                    maxScale = max(widthScale, heightScale)
                    scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                    # center the image if it's different size to the original
                    scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                    scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                    # now update
                    img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight

                if l.shape_id == 8 :
                    img = h.shapes[8]  # .element[0]
                    imgPic = img._pic
                    imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                    imgPart = h.part.related_parts[imgRID]

                    # get info about replacement image
                    with open(REPLACEMENT_IMG2, 'rb') as f:
                        rImgBlob = f.read()
                    rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
                    rImgWidth, rImgHeight = Pt(rImgWidth), Pt(rImgHeight)  # change from Px

                    # replace
                    imgPart._blob = rImgBlob

                    # now alter the size and position to suit that of the replacement img
                    # rescale sizes so image isn't stretched
                    widthScale = float(rImgWidth) / img.width
                    heightScale = float(rImgHeight) / img.height
                    maxScale = max(widthScale, heightScale)
                    scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
                    # center the image if it's different size to the original
                    scaledImgLeft = int(img.left + (img.width - scaledImgWidth) / 2)
                    scaledImgTop = int(img.top + (img.height - scaledImgHeight) / 2)
                    # now update
                    img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight



                #print(shape.shape_id)

            print("************************")

        prs.save("vvvvvvvv.pptx")
        p_.close()
    except :
        pass

def result():
    f = open('stor.bin', 'r')
    l = open('stor2.bin', 'r')
    m2 = load(l)
    m1 = load(f)
    # m = [[...............],[..........]]
    # m[0] == conttry 1 m[1] == contry 2 # m[0][-1] and m[1][-1]=== image
    contry1 = []
    contry2 = []

    for i in m1[:-1]:
        p=[]
        for e in i:
            p.append(transfr_mote(e))
        contry1.append(p)
        f.close()



    for h in m2[:-1]:
        L = []
        for mm in h:
            L.append(transfr_mote(mm))
        contry2.append(L)
        l.close()
    print contry1
    print '***************************'
    print contry2
    result_final([contry1,contry2])


def delet_(string):
    nomgre = ['0','1','2','3','4','5','6','7','8','9']
    s = ''
    for i in string :
        if i not in nomgre :
            s=s+''
        else :
            s = s + i
    if s != '':
        return  float(s)
    else :
        return 0


def transfr_mote(y):
    mote = [u'ترليون',u'مليون',u'الف',u'جانفي',u'فيفري',u'مارس',u'افريل',u'جوان',u'جويليا',u'اوت',u'أكتوبر',u'سبتمبر',u'نوفمبر',u'ديسمبر']
    if mote[0] in y :
        return  delet_(y)*10**12
    elif mote[1] in y :
        return  delet_(y)*10**9
    elif mote[2] in y :
        return  delet_(y)*10**3
    elif mote[3] in y or mote[4] in y or mote[5] in y or mote[6] in y or mote[7] in y or mote[8] in y or mote[9] in y or mote[10] in y or mote[11] in y or mote[12] in y or mote[13] in y :
        return 0
    else :
        return delet_(y)

def result_final(h):
    #[contry1,contry2]
    #print h[0]
    #print('*****************************')
    #print h[1]
    #print(h)
    #print h[1][-1][4]
    try:

        for i in h[0] :
            for e in i :
                if h[0][h[0].index(i)][i.index(e)] != h[0][1][0] and h[0][h[0].index(i)][i.index(e)] != h[0][-1][0] :
                    if e > h[1][h[0].index(i)][i.index(e)] :
                        #i[i.index(e)] = 'contry1'
                        h[0][h[0].index(i)][i.index(e)] = 'contry1'
                    elif e < h[1][h[0].index(i)][i.index(e)] :
                        #i[i.index(e)] = 'contry2'
                        h[0][h[0].index(i)][i.index(e)] = 'contry2'
                    else :
                        h[0][h[0].index(i)][i.index(e)] = 'nul'
                elif e < h[1][h[0].index(i)][i.index(e)] :
                    h[0][h[0].index(i)][i.index(e)] = 'contry1'
                elif e > h[1][h[0].index(i)][i.index(e)] :
                    h[0][h[0].index(i)][i.index(e)] = 'contry2'
                else :
                    h[0][h[0].index(i)][i.index(e)] = 'nul'

        #print('contry1 === ' ,h[0][1].count('contry1'))
        #print('contry2 ===', h[0][1].count('contry2'))
        #print('nul ===', h[0][1].count('nul'))
        p=[]
        for i in range(0,len(h[0])) :
            a = h[0][i].count('contry1')
            b = h[0][i].count('contry2')
            a_ = round((float(a*100))/(a+b),2)
            b_ = round((float(b*100))/(a+b),2)
            print '=======================' , a_
            p.append((a_,b_))
        cont1 = 0
        cont2 = 0
        for i in p :

            cont1 = i[0] + cont1
            cont2 = i[1] + cont2
        p.append((round(cont1/float(6),2) , round(cont2/float(6),2)))

        calcul(p)
    except :
        pass



result()