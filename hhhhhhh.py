from pptx import Presentation
from pptx.util import Pt
from PIL import Image

# inputs
PPT_TEMPLATE = 'C:/Users/benam/OneDrive/Bureau/cod python ppt/hemidi.pptx'
REPLACEMENT_IMG = 'hem.png'

# load preso and shape
presso = Presentation(PPT_TEMPLATE)
slide = presso.slides[0]
print len(slide.shapes)
"""
for k  in slide.shapes :
    print type(k)
"""

#print type(slide.shapes[41])


for g  in range(1 , 54):
    i = 1
    if i == 1 :
        #if g == 3 :

        slide = presso.slides[5]
        img = slide.shapes[g] #.element[0]
        m = slide.placeholders
        #print len(m)
        #print img.shape_id

        # get current image info
        if type(img) == type(slide.shapes[42]) :
            imgPic = img._pic
            imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
            imgPart = slide.part.related_parts[imgRID]

            # get info about replacement image
            with open(REPLACEMENT_IMG, 'rb') as f:
                rImgBlob = f.read()
            rImgWidth, rImgHeight = Image.open(REPLACEMENT_IMG).size
            rImgWidth, rImgHeight = Pt(rImgWidth), Pt(rImgHeight) # change from Px

            # replace
            imgPart._blob = rImgBlob

            # now alter the size and position to suit that of the replacement img
            # rescale sizes so image isn't stretched
            widthScale = float(rImgWidth) / img.width
            heightScale = float(rImgHeight) / img.height
            maxScale = max(widthScale, heightScale)
            scaledImgWidth, scaledImgHeight = int(rImgWidth / maxScale), int(rImgHeight / maxScale)
            # center the image if it's different size to the original
            scaledImgLeft = int(img.left + (img.width - scaledImgWidth)/2)
            scaledImgTop = int(img.top + (img.height - scaledImgHeight)/2)
            # now update
            img.left, img.top, img.width, img.height = scaledImgLeft, scaledImgTop, scaledImgWidth, scaledImgHeight
            presso.save("%s.pptx"%(g))
            print g , 'is photo '

        else :
            print '%s non photo'%(g)
print 'done'